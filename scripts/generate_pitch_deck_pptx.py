"""
Generate a pitch deck PPTX for Dream Flow from docs/pitch_deck/pitch_deck_slides.json.

Usage (Windows PowerShell):
  python scripts/generate_pitch_deck_pptx.py

Optional:
  python scripts/generate_pitch_deck_pptx.py --in docs/pitch_deck/pitch_deck_slides.json --out docs/pitch_deck/dream-flow-pitch-deck.pptx

Dependencies:
  pip install -r scripts/requirements-pitch-deck.txt
"""

from __future__ import annotations

import argparse
import json
from pathlib import Path
from typing import Any


def _require_python_pptx():
    try:
        from pptx import Presentation  # type: ignore
        from pptx.dml.color import RGBColor  # type: ignore
        from pptx.enum.shapes import MSO_SHAPE  # type: ignore
        from pptx.util import Inches, Pt  # type: ignore
        from pptx.enum.text import PP_ALIGN  # type: ignore
    except Exception as e:  # pragma: no cover
        raise SystemExit(
            "Missing dependency: python-pptx.\n"
            "Install it with:\n"
            "  pip install -r scripts/requirements-pitch-deck.txt\n"
            f"\nOriginal error: {e}"
        )
    return Presentation, Inches, Pt, PP_ALIGN, RGBColor, MSO_SHAPE


def _add_notes(slide, notes: str) -> None:
    if not notes:
        return
    notes_slide = slide.notes_slide
    text_frame = notes_slide.notes_text_frame
    text_frame.clear()
    text_frame.text = notes.strip()


def _set_footer(slide, footer: str, Inches, Pt, RGBColor):
    if not footer:
        return
    left = Inches(0.5)
    top = Inches(7.0)  # near bottom for 16:9
    width = Inches(12.3)
    height = Inches(0.3)
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = footer
    p.font.size = Pt(10)
    p.font.color.rgb = RGBColor(120, 120, 120)


def _add_title_slide(prs, s: dict[str, Any], meta: dict[str, Any], Inches, Pt, PP_ALIGN):
    layout = prs.slide_layouts[0]  # title
    slide = prs.slides.add_slide(layout)

    title = slide.shapes.title
    subtitle = slide.placeholders[1]

    title.text = s.get("title", meta.get("title", ""))
    subtitle.text = s.get("subtitle", meta.get("subtitle", ""))

    # Add bullets as an extra textbox (Gamma-like title slide)
    bullets = s.get("bullets") or []
    if bullets:
        left = Inches(1.1)
        top = Inches(3.3)
        width = Inches(11.0)
        height = Inches(3.0)
        box = slide.shapes.add_textbox(left, top, width, height)
        tf = box.text_frame
        tf.clear()
        for i, b in enumerate(bullets):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.text = str(b)
            p.level = 0
            p.font.size = Pt(20)
            p.space_after = Pt(6)

    _add_notes(slide, s.get("speaker_notes", ""))
    # RGBColor is provided via closure in build_pptx via partial application
    # (kept explicit in calls for clarity)
    _set_footer(slide, meta.get("footer", ""), Inches, Pt, meta["_RGBColor"])


def _add_bullets_slide(prs, s: dict[str, Any], meta: dict[str, Any], Inches, Pt):
    layout = prs.slide_layouts[1]  # title and content
    slide = prs.slides.add_slide(layout)

    slide.shapes.title.text = s.get("title", "")

    body = slide.shapes.placeholders[1].text_frame
    body.clear()

    subtitle = s.get("subtitle")
    if subtitle:
        p0 = body.paragraphs[0]
        p0.text = str(subtitle)
        p0.font.size = Pt(18)
        p0.space_after = Pt(10)
    else:
        body.paragraphs[0].text = ""

    bullets = s.get("bullets") or []
    for b in bullets:
        p = body.add_paragraph()
        p.text = str(b)
        p.level = 0
        p.font.size = Pt(20)
        p.space_after = Pt(4)

    _add_notes(slide, s.get("speaker_notes", ""))
    _set_footer(slide, meta.get("footer", ""), Inches, Pt, meta["_RGBColor"])


def _add_process_slide(prs, s: dict[str, Any], meta: dict[str, Any], Inches, Pt):
    # Use title-only layout and draw step boxes for a simple process diagram.
    layout = prs.slide_layouts[5]  # title only
    slide = prs.slides.add_slide(layout)
    slide.shapes.title.text = s.get("title", "")

    subtitle = s.get("subtitle", "")
    if subtitle:
        sub_box = slide.shapes.add_textbox(Inches(0.9), Inches(1.5), Inches(12.4), Inches(0.5))
        tf = sub_box.text_frame
        tf.text = subtitle
        tf.paragraphs[0].font.size = Pt(18)

    steps = s.get("steps") or []
    if steps:
        total = len(steps)
        left = 0.9
        top = 2.4
        gap = 0.25
        width = (13.3 - left * 2 - gap * (total - 1)) / max(total, 1)
        height = 2.1

        for i, step in enumerate(steps):
            x = Inches(left + i * (width + gap))
            y = Inches(top)
            w = Inches(width)
            h = Inches(height)
            shape = slide.shapes.add_shape(meta["_MSO_SHAPE"].ROUNDED_RECTANGLE, x, y, w, h)
            tf = shape.text_frame
            tf.clear()
            p = tf.paragraphs[0]
            p.text = f"{i+1}. {step}"
            p.font.size = Pt(18)

    _add_notes(slide, s.get("speaker_notes", ""))
    _set_footer(slide, meta.get("footer", ""), Inches, Pt, meta["_RGBColor"])


def build_pptx(deck: dict[str, Any], out_path: Path) -> None:
    Presentation, Inches, Pt, PP_ALIGN, RGBColor, MSO_SHAPE = _require_python_pptx()

    prs = Presentation()
    # 13.333 x 7.5 inches = standard 16:9 widescreen
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    meta: dict[str, Any] = {
        "title": deck.get("title", ""),
        "subtitle": deck.get("subtitle", ""),
        "footer": deck.get("footer", ""),
        "_RGBColor": RGBColor,
        "_MSO_SHAPE": MSO_SHAPE,
    }

    for s in deck.get("slides", []):
        t = s.get("type", "bullets")
        if t == "title":
            _add_title_slide(prs, s, meta, Inches, Pt, PP_ALIGN)
        elif t == "process":
            _add_process_slide(prs, s, meta, Inches, Pt)
        else:
            _add_bullets_slide(prs, s, meta, Inches, Pt)

    # Appendix
    for s in deck.get("appendix", []):
        _add_bullets_slide(prs, s, meta, Inches, Pt)

    out_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(out_path))


def main() -> int:
    parser = argparse.ArgumentParser()
    parser.add_argument(
        "--in",
        dest="in_path",
        default="docs/pitch_deck/pitch_deck_slides.json",
        help="Path to JSON slide spec",
    )
    parser.add_argument(
        "--out",
        dest="out_path",
        default="docs/pitch_deck/dream-flow-pitch-deck.pptx",
        help="Output PPTX path",
    )
    args = parser.parse_args()

    in_path = Path(args.in_path)
    out_path = Path(args.out_path)

    deck = json.loads(in_path.read_text(encoding="utf-8"))
    build_pptx(deck, out_path)
    print(f"Wrote: {out_path}")
    return 0


if __name__ == "__main__":
    raise SystemExit(main())

